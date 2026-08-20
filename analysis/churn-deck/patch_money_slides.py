# -*- coding: utf-8 -*-
"""Split the churn deck's money-case slide into two, and put controls on the prepay offer.

Reads  outputs/churn-strategy-2026-08-12.pptx   (14 slides, hand-built by a prior session)
Writes outputs/churn-strategy-2026-08-13.pptx   (15 slides)

Slide 11 of the source deck carried both money cases at once: the AM programme
break-even (a cost-to-required-reduction division, drawn as a 4-bar chart that
plots a straight line through the origin) and the prepay P&L (four rows whose
three governing assumptions appeared nowhere on the slide). The arithmetic was
right to the rupiah; the presentation buried it. This script rebuilds slide 11
as Decision 1's money case, inserts a new slide 12 for Decision 2's, adds the
prepay authorisation line to slide 9, and fixes the two forward slide references
that break when a slide is inserted at position 12.

Every figure that reaches a slide is recomputed here from either the target-list
workbook or the documented model inputs, and asserted against the number the
prior deck showed. Nothing is typed as a display literal.

Run:  python analysis/churn-deck/patch_money_slides.py
"""

from __future__ import annotations

import copy
import shutil
import sys
import tempfile
from pathlib import Path

import openpyxl
from pptx import Presentation
from pptx.oxml.ns import qn
from lxml import etree

REPO = Path(__file__).resolve().parents[2]
SRC_DECK = REPO / "outputs" / "churn-strategy-2026-08-12.pptx"
OUT_DECK = REPO / "outputs" / "churn-strategy-2026-08-13.pptx"
TARGET_LIST = REPO / "outputs" / "churn-am-target-list-2026-08-11.xlsx"


# --------------------------------------------------------------------------
# 1. The model. Recomputed, then checked against what the prior deck displayed.
# --------------------------------------------------------------------------

# Inputs that come from the v2 memo rather than from a machine-readable tab.
# churn-am-program-v2-2026-08-12.md: §1.1 for the monthly book, §2.1 for the
# forward hazard, §3.7 for L and the cost of capital.
MONTHLY_BOOK_ARR = 90.79e9   # §1.1  active ARR billed monthly (79.5% of the book)
CONVERSION_SHARE = 0.10      #       illustrative share of that book converted
DISCOUNT = 0.08              # §3.7  the rate we are asking to approve
COST_OF_CAPITAL = 0.11       # §3.7  Finance owns this
PREPAY_MONTHS_EARLY = 6      #       annual up front vs monthly in arrears
W_FORWARD_12M = 0.129        # §2.1  P(churn in next 12 mo) for an account at month 12
L_RECOVERED = 0.60           # §3.7  share of a failed account's annual value recovered
ADVERSE_HAIRCUT = 0.50       # §3.7  healthiest accounts self-select in first

# Programme-cost scenarios shown on Decision 1's slide.
COST_SCENARIOS = (0.5e9, 1.0e9, 1.5e9, 2.0e9)
ANCHOR_COST = 1.0e9

# -- AM coverage ------------------------------------------------------------
# The accounts two AMs already own, as named by Grady on 13 Aug 2026 and
# resolved to exact target-list records. Three needed manual resolution:
# "Brinks" fuzzy-matched to "BM Trans" at 0.57 but is Brinks Indonesia (rank 2,
# IDR 4.90B); "Sanobar" and "Duta energi" each had two candidates in the
# register and Grady confirmed the ones below.
COVERED_ACCOUNTS = (
    "Patra Logistik, PT",
    "Brinks Indonesia, PT",
    "IndoMobil Bussan Trucking, PT",
    "Bank Syariah Indonesia, PT",
    "Swadharma Sarana Informatika, PT",
    "Wastec International",
    "Sany Makmur Perkasa, PT",
    "Jasa Berdikari Logistics, PT",
    '"K" Line Mobaru Diamond Indonesia, PT',
    "Serikat Hantar Ekspedisi, PT (Seryu Cargo)",
    "Tempirai Energy Resources, PT",
    "Farika Beton, PT",
    "Almera Sukses Sejahtera, PT",
    "Adarton Indo Pacific, PT",
    "Berkat Jaya Beton, PT",
    "Karya Semesta Logistik, PT",
    "Meindo Elang Indah, PT",
    "Skuad Manajemen Indonesia, PT",
    "Hijrah Gizi Hewani, PT (Hijrah Food)",
    "Widya Sapta Contractor, PT",
    "Andalan Artha Primanusa, PT",
    "Indosarana Jaya Perkasa, PT",
    "Tata Bara Utama, PT",
    "Cipta Unggul Lintas Samudra, PT",
    "IMC Transporindologistik, PT",
    "Putera Persada Jaya, PT",
    "Triputra Menara Jaya, PT",
    "Jangkar Pasifik Transport, PT",
    "Antero Bahana Cemerlang (ABC Express), PT",
    "Sanobar Gunajaya, PT",
    "Jidousha Niaga Logistik Jnl, PT",
    "Gobel Dharma Sarana Karya (GDSK), PT",
    "Duta Energi Muliatama, PT",
)

AM_SALARY_MONTH = 12e6      # per AM
AM_INCENTIVE_MONTHS_PER_Q = 1
AM_ENGAGEMENT_MONTH = 5e6   # customer engagement, per AM
AMS_TODAY = 2
AMS_TO_ADD = 2              # Grady's call: 2 more cover the remaining top-50
SMALL_ACCOUNT_FLOOR = 100e6  # below this, an account belongs in Play C


def read_pool_from_target_list() -> dict:
    """Recompute the covered pool from the AM target list, excluding Corin.

    OneDrive holds a lock on files in outputs/ often enough that direct reads
    raise PermissionError; copy to a temp path first.
    """
    tmp = Path(tempfile.gettempdir()) / "_churn_target_list.xlsx"
    shutil.copy2(TARGET_LIST, tmp)
    wb = openpyxl.load_workbook(tmp, data_only=True)
    ws = wb["At_Risk_Active_Accounts"]

    rows = []
    for row in ws.iter_rows(min_row=5, values_only=True):
        name, arr, loss = row[0], row[4], row[10]
        if name and isinstance(arr, (int, float)) and isinstance(loss, (int, float)):
            rows.append((name, float(arr), float(loss)))

    ex_corin = [r for r in rows if "Corin" not in r[0]]
    ex_corin.sort(key=lambda r: -r[2])
    top50 = ex_corin[:50]

    return {
        "accounts": len(ex_corin),
        "total_arr": sum(r[1] for r in ex_corin),
        "total_loss": sum(r[2] for r in ex_corin),
        "top50_arr": sum(r[1] for r in top50),
        "top50_loss": sum(r[2] for r in top50),
        "rows": ex_corin,
        "top50": top50,
    }


def build_coverage(pool: dict) -> dict:
    """What the existing AM team already owns, and what two more would cost.

    Every covered name must resolve to a target-list record: a renamed or
    dropped record should fail the build, not silently shrink coverage.
    """
    by_name = {r[0]: r for r in pool["rows"]}
    missing = [n for n in COVERED_ACCOUNTS if n not in by_name]
    if missing:
        raise SystemExit("covered accounts not found in target list: %r" % missing)

    covered = [by_name[n] for n in COVERED_ACCOUNTS]
    top50_names = {r[0] for r in pool["top50"]}
    cov_in_50 = [r for r in covered if r[0] in top50_names]
    uncovered_50 = [r for r in pool["top50"] if r[0] not in set(COVERED_ACCOUNTS)]
    small = [r for r in covered if r[1] < SMALL_ACCOUNT_FLOOR]

    per_am = (AM_SALARY_MONTH * 12
              + AM_SALARY_MONTH * AM_INCENTIVE_MONTHS_PER_Q * 4
              + AM_ENGAGEMENT_MONTH * 12)
    committed = per_am * AMS_TODAY
    incremental = per_am * AMS_TO_ADD
    ams_total = AMS_TODAY + AMS_TO_ADD

    unc_loss = sum(r[2] for r in uncovered_50)
    return {
        "covered_n": len(covered),
        "covered_arr": sum(r[1] for r in covered),
        "covered_loss": sum(r[2] for r in covered),
        "in50_n": len(cov_in_50),
        "in50_arr": sum(r[1] for r in cov_in_50),
        "in50_loss": sum(r[2] for r in cov_in_50),
        "in50_share": sum(r[2] for r in cov_in_50) / pool["top50_loss"],
        "unc_n": len(uncovered_50),
        "unc_arr": sum(r[1] for r in uncovered_50),
        "unc_loss": unc_loss,
        "outside_n": len(covered) - len(cov_in_50),
        "outside_arr": sum(r[1] for r in covered) - sum(r[1] for r in cov_in_50),
        "outside_loss": sum(r[2] for r in covered) - sum(r[2] for r in cov_in_50),
        "small_n": len(small),
        "small_loss": sum(r[2] for r in small),
        "per_am": per_am,
        "salary": AM_SALARY_MONTH * 12,
        "incentive": AM_SALARY_MONTH * AM_INCENTIVE_MONTHS_PER_Q * 4,
        "engagement": AM_ENGAGEMENT_MONTH * 12,
        "committed": committed,
        "incremental": incremental,
        "programme": committed + incremental,
        "ams_total": ams_total,
        # The bar the incremental ask has to clear: it buys coverage of the
        # accounts nobody owns, so it is measured against that pool alone.
        "bar_incremental": incremental / unc_loss,
        "bar_programme": (committed + incremental) / pool["top50_loss"],
        "load": (len(covered) + len(uncovered_50)) / ams_total,
        "load_ex_small": (len(covered) - len(small) + len(uncovered_50)) / ams_total,
    }


def build_model() -> dict:
    pool = read_pool_from_target_list()
    covered = pool["top50_loss"]

    # Decision 1 -- one division. Required RELATIVE churn reduction on the
    # covered pool for an annual programme cost to pay for itself in year 1.
    ladder = [(c, c / covered) for c in COST_SCENARIOS]

    # Decision 2 -- the prepay, per rupiah converted and at the illustration.
    converted = MONTHLY_BOOK_ARR * CONVERSION_SHARE
    early_frac = PREPAY_MONTHS_EARLY / 12.0

    def prepay(base: float) -> dict:
        discount = -base * DISCOUNT
        working_capital = base * early_frac
        carry = working_capital * COST_OF_CAPITAL
        retained = base * W_FORWARD_12M * L_RECOVERED * ADVERSE_HAIRCUT
        return {
            "discount": discount,
            "carry": carry,
            "retained": retained,
            "net": discount + carry + retained,
            "working_capital": working_capital,
        }

    # Break-even discount is the same model rearranged: the value the lock buys,
    # plus the carry on cash pulled forward, expressed as a share of the base.
    be_average = (W_FORWARD_12M * L_RECOVERED) + (early_frac * COST_OF_CAPITAL)
    be_adverse = (W_FORWARD_12M * ADVERSE_HAIRCUT * L_RECOVERED) + (early_frac * COST_OF_CAPITAL)

    return {
        "pool": pool,
        "cov": build_coverage(pool),
        "covered": covered,
        "ladder": ladder,
        "anchor_pct": ANCHOR_COST / covered,
        "converted": converted,
        "unit": prepay(1.0e9),
        "illus": prepay(converted),
        "be_average": be_average,
        "be_adverse": be_adverse,
    }


def check(model: dict) -> None:
    """Assert the recomputed model reproduces every figure the prior deck showed.

    If an upstream extract changes, this fails loudly instead of letting the
    slide drift away from the analysis it claims to rest on.
    """
    p, u, i = model["pool"], model["unit"], model["illus"]
    cases = [
        ("accounts ex-Corin", p["accounts"], 2002, 0),
        ("total expected loss (B)", p["total_loss"] / 1e9, 13.23, 0.01),
        ("total active ARR (B)", p["total_arr"] / 1e9, 94.99, 0.01),
        ("top-50 ARR (B)", p["top50_arr"] / 1e9, 50.95, 0.01),
        ("top-50 loss (B)", model["covered"] / 1e9, 8.70, 0.01),
        ("top-50 share of loss (%)", 100 * model["covered"] / p["total_loss"], 65.8, 0.1),
        ("anchor reduction (%)", 100 * model["anchor_pct"], 11.5, 0.05),
        ("converted ARR (B)", model["converted"] / 1e9, 9.08, 0.01),
        ("unit discount (M)", u["discount"] / 1e6, -80, 0.5),
        ("unit carry (M)", u["carry"] / 1e6, 55, 0.5),
        ("unit retained (M)", u["retained"] / 1e6, 39, 0.5),
        ("unit net (M)", u["net"] / 1e6, 14, 0.5),
        ("unit working capital (M)", u["working_capital"] / 1e6, 500, 0.5),
        ("illus discount (M)", i["discount"] / 1e6, -726, 1),
        ("illus carry (M)", i["carry"] / 1e6, 499, 1),
        ("illus retained (M)", i["retained"] / 1e6, 352, 1),
        ("illus net (M)", i["net"] / 1e6, 124, 1),
        ("illus working capital (B)", i["working_capital"] / 1e9, 4.54, 0.01),
        ("break-even, average (%)", 100 * model["be_average"], 13.2, 0.05),
        ("break-even, adverse (%)", 100 * model["be_adverse"], 9.4, 0.05),
    ]
    c = model["cov"]
    cases += [
        ("AM-covered accounts", c["covered_n"], 33, 0),
        ("AM-covered ARR (B)", c["covered_arr"] / 1e9, 25.26, 0.01),
        ("AM-covered loss (B)", c["covered_loss"] / 1e9, 4.34, 0.01),
        ("covered in top 50", c["in50_n"], 13, 0),
        ("covered in top 50, ARR (B)", c["in50_arr"] / 1e9, 21.93, 0.01),
        ("covered in top 50, loss (B)", c["in50_loss"] / 1e9, 4.03, 0.01),
        ("share of top-50 loss (%)", 100 * c["in50_share"], 46.3, 0.1),
        ("top 50 with no owner", c["unc_n"], 37, 0),
        ("no-owner ARR (B)", c["unc_arr"] / 1e9, 29.02, 0.01),
        ("no-owner loss (B)", c["unc_loss"] / 1e9, 4.67, 0.01),
        ("sub-100M covered accts", c["small_n"], 10, 0),
        ("cost per AM (M)", c["per_am"] / 1e6, 252, 0.5),
        ("committed, 2 AMs (M)", c["committed"] / 1e6, 504, 0.5),
        ("incremental, 2 AMs (M)", c["incremental"] / 1e6, 504, 0.5),
        ("programme, 4 AMs (M)", c["programme"] / 1e6, 1008, 0.5),
        ("incremental bar (%)", 100 * c["bar_incremental"], 10.8, 0.05),
        ("programme bar (%)", 100 * c["bar_programme"], 11.6, 0.05),
        ("accounts per AM", c["load"], 17.5, 0.05),
        ("accounts per AM ex-small", c["load_ex_small"], 15.0, 0.05),
    ]
    bad = [(n, got, want) for n, got, want, tol in cases if abs(got - want) > tol]
    for name, got, want in bad:
        print(f"  MISMATCH {name}: recomputed {got:.4f}, deck showed {want}", file=sys.stderr)
    if bad:
        raise SystemExit("model no longer reproduces the deck's figures -- refusing to build")
    for name, got, want, _ in cases:
        print(f"  ok  {name:28s} {got:>12.3f}")


# --------------------------------------------------------------------------
# 2. Design tokens, lifted from the source deck's own XML.
# --------------------------------------------------------------------------

BG_LIGHT = "F5F7FA"
NAVY = "16264A"
INK = "1B2437"
MUTED = "6B7280"
TEAL = "13908C"
GOLD = "E39B2E"
RED = "D9534F"
PALE_TEAL = "EDF5F4"
ON_DARK = "FFFFFF"
ON_DARK_SOFT = "CADCFC"
HAIRLINE = "333333"

TITLE_FONT = "Cambria"
BODY_FONT = "Calibri"

EMU_PER_IN = 914400
SHADOW = (
    '<a:effectLst><a:outerShdw sx="100000" sy="100000" kx="0" ky="0" algn="bl"'
    ' rotWithShape="0" blurRad="127000" dist="25400" dir="5400000">'
    f'<a:srgbClr val="9AA5B4"><a:alpha val="28000"/></a:srgbClr>'
    "</a:outerShdw></a:effectLst>"
)
NSMAP = (
    'xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
    'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
    'xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"'
)


def emu(inches: float) -> int:
    return int(round(inches * EMU_PER_IN))


def esc(text: str) -> str:
    return text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")


class Canvas:
    """Emits shapes as raw XML so they match the hand-built slides exactly."""

    def __init__(self, slide, next_id: int = 2):
        self.slide = slide
        self.tree = slide.shapes._spTree
        self.next_id = next_id

    def _add(self, xml: str) -> None:
        self.tree.append(etree.fromstring(xml))
        self.next_id += 1

    def panel(self, x, y, w, h, fill, *, adj=None, line=HAIRLINE, line_w=12700,
              shadow=False, ellipse=False):
        i = self.next_id
        if ellipse:
            geom = '<a:prstGeom prst="ellipse"><a:avLst/></a:prstGeom>'
        else:
            inner = f'<a:gd name="adj" fmla="val {adj}"/>' if adj is not None else ""
            geom = f'<a:prstGeom prst="roundRect"><a:avLst>{inner}</a:avLst></a:prstGeom>'
        fill_xml = f'<a:solidFill><a:srgbClr val="{fill}"/></a:solidFill>' if fill else "<a:noFill/>"
        line_xml = (
            f'<a:ln w="{line_w}"><a:solidFill><a:srgbClr val="{line}"/></a:solidFill>'
            '<a:prstDash val="solid"/></a:ln>'
        ) if line else "<a:ln/>"
        self._add(
            f'<p:sp {NSMAP}><p:nvSpPr><p:cNvPr id="{i}" name="Shape {i - 1}"/>'
            "<p:cNvSpPr/><p:nvPr/></p:nvSpPr><p:spPr>"
            f'<a:xfrm><a:off x="{emu(x)}" y="{emu(y)}"/>'
            f'<a:ext cx="{emu(w)}" cy="{emu(h)}"/></a:xfrm>'
            f"{geom}{fill_xml}{line_xml}{SHADOW if shadow else ''}</p:spPr></p:sp>"
        )

    def text(self, x, y, w, h, runs, *, align="l", line_pts=None, anchor="ctr"):
        """runs: sequence of (text, size_pt, bold, colour, font)."""
        i = self.next_id
        lnspc = f'<a:lnSpc><a:spcPts val="{int(line_pts * 100)}"/></a:lnSpc>' if line_pts else ""
        algn = f' algn="{align}"' if align != "l" else ""
        ppr = f'<a:pPr{algn} indent="0" marL="0">{lnspc}<a:buNone/></a:pPr>'
        body = [ppr]
        last_sz = 1100
        for text, sz, bold, colour, font in runs:
            last_sz = int(sz * 100)
            b = ' b="1"' if bold else ""
            body.append(
                f'<a:r><a:rPr lang="en-US" sz="{last_sz}"{b} dirty="0">'
                f'<a:solidFill><a:srgbClr val="{colour}"/></a:solidFill>'
                f'<a:latin typeface="{font}" pitchFamily="34" charset="0"/>'
                f'<a:ea typeface="{font}" pitchFamily="34" charset="-122"/>'
                f'<a:cs typeface="{font}" pitchFamily="34" charset="-120"/>'
                f"</a:rPr><a:t>{esc(text)}</a:t></a:r>"
            )
        body.append(f'<a:endParaRPr lang="en-US" sz="{last_sz}" dirty="0"/>')
        self._add(
            f'<p:sp {NSMAP}><p:nvSpPr><p:cNvPr id="{i}" name="Text {i - 1}"/>'
            "<p:cNvSpPr/><p:nvPr/></p:nvSpPr><p:spPr>"
            f'<a:xfrm><a:off x="{emu(x)}" y="{emu(y)}"/>'
            f'<a:ext cx="{emu(w)}" cy="{emu(h)}"/></a:xfrm>'
            '<a:prstGeom prst="rect"><a:avLst/></a:prstGeom><a:noFill/><a:ln/></p:spPr>'
            '<p:txBody><a:bodyPr wrap="square" lIns="0" tIns="0" rIns="0" bIns="0"'
            f' rtlCol="0" anchor="{anchor}"/><a:lstStyle/><a:p>{"".join(body)}</a:p>'
            "</p:txBody></p:sp>"
        )

    def text_stack(self, x, y, w, h, paras, *, align="l", anchor="ctr"):
        """Several paragraphs in one box. paras: [(runs, line_pts), ...]."""
        i = self.next_id
        algn = f' algn="{align}"' if align != "l" else ""
        body = []
        last_sz = 1100
        for runs, line_pts in paras:
            lnspc = f'<a:lnSpc><a:spcPts val="{int(line_pts * 100)}"/></a:lnSpc>'
            body.append(f"<a:p><a:pPr{algn} indent=\"0\" marL=\"0\">{lnspc}<a:buNone/></a:pPr>")
            for text, sz, bold, colour, font in runs:
                last_sz = int(sz * 100)
                b = ' b="1"' if bold else ""
                body.append(
                    f'<a:r><a:rPr lang="en-US" sz="{last_sz}"{b} dirty="0">'
                    f'<a:solidFill><a:srgbClr val="{colour}"/></a:solidFill>'
                    f'<a:latin typeface="{font}" pitchFamily="34" charset="0"/>'
                    f'<a:ea typeface="{font}" pitchFamily="34" charset="-122"/>'
                    f'<a:cs typeface="{font}" pitchFamily="34" charset="-120"/>'
                    f"</a:rPr><a:t>{esc(text)}</a:t></a:r>"
                )
            body.append(f'<a:endParaRPr lang="en-US" sz="{last_sz}" dirty="0"/></a:p>')
        self._add(
            f'<p:sp {NSMAP}><p:nvSpPr><p:cNvPr id="{i}" name="Text {i - 1}"/>'
            "<p:cNvSpPr/><p:nvPr/></p:nvSpPr><p:spPr>"
            f'<a:xfrm><a:off x="{emu(x)}" y="{emu(y)}"/>'
            f'<a:ext cx="{emu(w)}" cy="{emu(h)}"/></a:xfrm>'
            '<a:prstGeom prst="rect"><a:avLst/></a:prstGeom><a:noFill/><a:ln/></p:spPr>'
            '<p:txBody><a:bodyPr wrap="square" lIns="0" tIns="0" rIns="0" bIns="0"'
            f' rtlCol="0" anchor="{anchor}"/><a:lstStyle/>{"".join(body)}'
            "</p:txBody></p:sp>"
        )

    # -- composite helpers -------------------------------------------------

    def badge(self, x, y, size, letter, fill):
        self.panel(x, y, size, size, fill, ellipse=True)
        self.text(x, y, size, size, [(letter, 15, True, ON_DARK, BODY_FONT)], align="ctr")

    def eyebrow(self, label, colour=TEAL):
        self.text(0.6, 0.34, 12.1, 0.28, [(label, 11, True, colour, BODY_FONT)])

    def title(self, label, colour=INK):
        self.text(0.6, 0.64, 12.1, 0.82, [(label, 30, True, colour, TITLE_FONT)])

    def banner(self, lead, rest):
        self.panel(0.6, 1.6, 12.1, 0.52, NAVY, adj=13462)
        self.text(0.96, 1.6, 11.38, 0.52, [
            (lead + "  ", 12.5, True, GOLD, BODY_FONT),
            (rest, 12.5, False, ON_DARK_SOFT, BODY_FONT),
        ])

    def card(self, x, y, w, h):
        self.panel(x, y, w, h, ON_DARK, adj=2646, shadow=True)

    def callout(self, runs):
        self.panel(0.6, 6.28, 12.1, 0.78, ON_DARK, adj=11538, line=TEAL, line_w=15875, shadow=True)
        self.badge(0.86, 6.51, 0.32, "!", TEAL)
        self.text(1.32, 6.4, 11.1, 0.56, runs, line_pts=17)

    def row(self, y, label_x, label_w, label, cols, *, bold=False, h=0.38,
            sublabel=None):
        if sublabel:
            # Two stacked lines, top-anchored, so the value still centres on
            # the block. Used where a row's arithmetic is not self-evident.
            self.text_stack(label_x, y, label_w, h + 0.20, [
                ([(label, 11.5, bold, INK, BODY_FONT)], 15),
                ([(sublabel, 9.5, False, MUTED, BODY_FONT)], 12),
            ], anchor="t")
        else:
            self.text(label_x, y, label_w, h,
                      [(label, 11.5, bold, INK, BODY_FONT)], line_pts=15)
        for cx, cw, value, colour in cols:
            self.text(cx, y, cw, h, [(value, 14, True, colour, TITLE_FONT)], align="r")


def wipe(slide) -> None:
    """Remove every shape from a slide, dropping any chart relationships."""
    tree = slide.shapes._spTree
    for el in list(tree):
        if el.tag in (qn("p:nvGrpSpPr"), qn("p:grpSpPr")):
            continue
        for ref in el.iter():
            rid = ref.get(qn("r:id")) or ref.get(qn("r:embed"))
            if rid and rid in slide.part.rels:
                slide.part.drop_rel(rid)
        tree.remove(el)


def set_notes(slide, text: str) -> None:
    slide.notes_slide.notes_text_frame.text = text


def money(v: float, unit: str) -> str:
    """Signed, rounded figure with the deck's U+2212 minus sign."""
    scale = 1e6 if unit == "M" else 1e9
    n = v / scale
    sign = "−" if n < 0 else "+"
    if unit == "M":
        return f"{sign}{abs(n):,.0f}M"
    return f"{sign}{abs(n):,.2f}B"


# --------------------------------------------------------------------------
# 3. Coverage -- what the existing AM team already owns.
# --------------------------------------------------------------------------

def build_coverage_slide(slide, m: dict) -> None:
    c = Canvas(slide)
    v = m["cov"]

    c.eyebrow("COVERAGE")
    c.title(f"We already cover {v['in50_share'] * 100:.0f}% of the top-50 exposure")
    c.banner(
        "This is not a new team.",
        f"{AMS_TODAY} AMs cover {v['in50_n']} of the top 50 today. The ask is "
        f"{AMS_TO_ADD} more, to finish the other {v['unc_n']}.",
    )

    # -- left card: the coverage gap ---------------------------------------
    c.card(0.6, 2.3, 6.02, 3.78)
    c.badge(0.88, 2.5, 0.42, "✓", TEAL)
    c.text(1.44, 2.49, 5.0, 0.32, [("Covered today", 17, True, INK, TITLE_FONT)])
    c.text(0.9, 2.92, 5.42, 0.32, [(
        f"{AMS_TODAY} AMs, {v['covered_n']} accounts, IDR "
        f"{v['committed'] / 1e6:,.0f}M a year.", 11.5, False, MUTED, BODY_FONT)],
        line_pts=15)

    lx, lw = 1.08, 2.0
    ax, aw = 3.10, 0.80
    rx, rw = 4.00, 1.10
    ex, ew = 5.20, 1.02
    for cx, cw, head in ((ax, aw, "ACCTS"), (rx, rw, "ARR"), (ex, ew, "EXPOSURE")):
        c.text(cx, 3.40, cw, 0.24, [(head, 9.5, True, MUTED, BODY_FONT)], align="r")

    table = [
        (3.72, "Top 50, total", 50, m["pool"]["top50_arr"], m["pool"]["top50_loss"],
         INK, False),
        (4.20, f"Covered by {AMS_TODAY} AMs", v["in50_n"], v["in50_arr"],
         v["in50_loss"], TEAL, True),
        (4.68, "No owner today", v["unc_n"], v["unc_arr"], v["unc_loss"], RED, True),
    ]
    c.panel(0.9, 4.16, 5.42, 0.46, PALE_TEAL, adj=10000)
    for y, label, n, arr, loss, colour, bold in table:
        c.text(lx, y, lw, 0.38, [(label, 11.5, bold, INK, BODY_FONT)], line_pts=15)
        c.text(ax, y, aw, 0.38, [(f"{n}", 14, True, colour, TITLE_FONT)], align="r")
        c.text(rx, y, rw, 0.38,
               [(f"{arr / 1e9:.2f}B", 14, True, colour, TITLE_FONT)], align="r")
        c.text(ex, y, ew, 0.38,
               [(f"{loss / 1e9:.2f}B", 14, True, colour, TITLE_FONT)], align="r")

    c.text(1.08, 5.28, 5.14, 0.62, [(
        f"The other {v['outside_n']} covered accounts sit outside the top 50 — "
        f"IDR {v['outside_arr'] / 1e9:.2f}B of ARR but only "
        f"{v['outside_loss'] / 1e9:.2f}B of exposure. {v['small_n']} are under "
        f"IDR {SMALL_ACCOUNT_FLOOR / 1e6:.0f}M and belong in Play C.",
        10.5, False, MUTED, BODY_FONT)], line_pts=13, anchor="t")

    # -- right card: what two more cost -----------------------------------
    c.card(7.28, 2.3, 5.42, 3.78)
    c.badge(7.56, 2.5, 0.42, "+", GOLD)
    c.text(8.12, 2.49, 4.4, 0.32,
           [(f"What {AMS_TO_ADD} more AMs cost", 17, True, INK, TITLE_FONT)])
    c.text(7.58, 2.92, 4.84, 0.32, [(
        f"IDR {v['per_am'] / 1e6:.0f}M per AM = {v['salary'] / 1e6:.0f}M salary + "
        f"{v['incentive'] / 1e6:.0f}M incentive + {v['engagement'] / 1e6:.0f}M "
        "engagement.", 10.5, False, MUTED, BODY_FONT)], line_pts=13, anchor="ctr")

    klx, kvx, kvw = 7.62, 10.60, 1.80
    cost_rows = [
        (3.44, "Per AM, per year", v["per_am"], INK, False),
        (3.88, f"{AMS_TODAY} AMs today — already committed", v["committed"],
         MUTED, False),
        (4.32, f"+ {AMS_TO_ADD} more AMs — the ask", v["incremental"], GOLD, True),
        (4.76, f"Total programme, {v['ams_total']} AMs", v["programme"], NAVY, True),
    ]
    c.panel(7.5, 4.72, 5.0, 0.46, PALE_TEAL, adj=10000)
    for y, label, amount, colour, bold in cost_rows:
        c.text(klx, y, 2.9, 0.38, [(label, 11.5, bold, INK, BODY_FONT)], line_pts=15)
        c.text(kvx, y, kvw, 0.38,
               [(f"{amount / 1e6:,.0f}M", 14, True, colour, TITLE_FONT)], align="r")

    c.panel(7.5, 5.30, 5.0, 0.62, NAVY, adj=8000)
    c.text(7.68, 5.38, 4.64, 0.46, [
        (f"{v['incremental'] / 1e6:,.0f}M against the "
         f"{v['unc_loss'] / 1e9:.2f}B nobody owns", 11.5, True, ON_DARK, BODY_FONT),
        (f" = a {v['bar_incremental'] * 100:.1f}% relative churn reduction to "
         "break even.", 11.5, False, ON_DARK_SOFT, BODY_FONT),
    ], line_pts=14)

    c.callout([
        (f"The ask is {v['incremental'] / 1e6:,.0f}M, not "
         f"{v['programme'] / 1e6:,.0f}M", 13.5, True, INK, BODY_FONT),
        (f" — half the programme is already funded. Each of the {v['ams_total']} "
         f"then carries about {v['load']:.0f} accounts, or {v['load_ex_small']:.0f} "
         f"if the {v['small_n']} sub-{SMALL_ACCOUNT_FLOOR / 1e6:.0f}M accounts move "
         "to Play C. That is the assumption to test first.",
         13.5, False, INK, BODY_FONT),
    ])

    set_notes(slide, (
        "This slide exists because the earlier draft claimed the top 50 have no owner. They do -- "
        "13 of them, carrying 46% of the expected loss. Better to correct it myself than have it "
        "corrected in the room.\n"
        "The reframe matters: this is not 'fund a retention team', it is 'finish a job that is "
        "already half done and already showing where the exposure sits'.\n"
        "The number I am least sure of is the load. Two more AMs means roughly 17 accounts each, "
        "and the 37 uncovered are heavier than what we carry now -- every one is above IDR 100M and "
        "19 are above 500M, where today we carry 10 accounts under 100M. If it proves too thin the "
        "fix is to move those 10 to Play C first, not to hire a fifth AM.\n"
        "One data caveat if pressed on the exposure figures: Jasa Berdikari and K Line Mobaru show "
        "zero expected loss despite IDR 1.30B of ARR between them, because their tenure sits past "
        "the point the hazard model can measure. Book-wide that affects 174 accounts holding 10.43B. "
        "Zero modelled risk is not zero risk."
    ))


# --------------------------------------------------------------------------
# 4. Slide 11 -- Decision 1's money case.
# --------------------------------------------------------------------------

def build_decision_1(slide, m: dict) -> None:
    wipe(slide)
    c = Canvas(slide)

    c.eyebrow("THE MONEY CASE · DECISION 1")
    c.title("What the AM programme has to earn")
    c.banner(
        "No uplift is quoted.",
        "I have never run this play, so instead of a forecast: here is the bar it has "
        "to clear, and exactly what is and is not counted in it.",
    )

    # -- left card: the division ------------------------------------------
    c.card(0.6, 2.3, 6.02, 3.78)
    c.badge(0.88, 2.5, 0.42, "A", TEAL)
    c.text(1.44, 2.49, 5.0, 0.32,
           [("Programme cost — the AM team", 17, True, INK, TITLE_FONT)])
    # One line only -- two lines here collide with the formula strip below.
    c.text(0.9, 2.92, 5.42, 0.32, [(
        f"Covered pool: top 50 accounts — IDR {m['covered'] / 1e9:.2f}B of the "
        f"IDR {m['pool']['total_loss'] / 1e9:.2f}B expected loss.",
        11.5, False, MUTED, BODY_FONT)], line_pts=15)

    c.panel(0.9, 3.34, 5.42, 0.42, PALE_TEAL, adj=11364)
    c.text(0.9, 3.34, 5.42, 0.42, [
        ("relative churn reduction needed  =  total annual cost  ÷  IDR "
         f"{m['covered'] / 1e9:.2f}B", 11.5, True, INK, BODY_FONT),
    ], align="ctr")

    # Inset from the 0.9..6.32 highlight panel so values are not clipped by it.
    label_x, label_w = 1.08, 2.9
    val_x, val_w = 3.98, 2.16
    c.text(label_x, 3.94, label_w, 0.26,
           [("TOTAL PROGRAMME COST", 10, True, MUTED, BODY_FONT)])
    c.text(val_x, 3.94, val_w, 0.26,
           [("REDUCTION NEEDED", 10, True, MUTED, BODY_FONT)], align="r")

    for idx, (cost, pct) in enumerate(m["ladder"]):
        y = 4.24 + idx * 0.44
        if abs(cost - ANCHOR_COST) < 1:
            c.panel(0.9, y - 0.02, 5.42, 0.42, PALE_TEAL, adj=11364)
        bold = abs(cost - ANCHOR_COST) < 1
        c.row(y, label_x, label_w, f"IDR {cost / 1e9:.1f}B",
              [(val_x, val_w, f"{pct * 100:.1f}%", INK)], bold=bold)

    # -- right card: what the number does and does not count --------------
    c.card(7.28, 2.3, 5.42, 3.78)
    c.text(7.56, 2.49, 4.9, 0.32,
           [(f"What the {m['anchor_pct'] * 100:.1f}% counts", 17, True, INK, TITLE_FONT)])
    c.text(7.56, 2.88, 4.86, 0.28, [
        ("A threshold, not a forecast. Here is the accounting behind it.",
         11.5, False, MUTED, BODY_FONT)], line_pts=15)

    blocks = [
        (3.32, "COUNTED", TEAL,
         "ARR retained in year 1, at 100% of revenue.", 0.34),
        (4.00, "LEFT OUT — MAKES THE BAR LOOK EASIER", GOLD,
         "No gross-margin haircut. A saved rupiah of ARR is revenue, not profit.", 0.48),
        (4.82, "LEFT OUT — MAKES THE BAR LOOK HARDER", GOLD,
         "An account saved in year 1 keeps paying in year 2 and beyond. "
         "Only the first year is credited.", 0.48),
    ]
    for y, label, colour, body, bh in blocks:
        c.text(7.56, y, 4.86, 0.22, [(label, 10, True, colour, BODY_FONT)])
        c.text(7.56, y + 0.22, 4.86, bh, [(body, 11, False, INK, BODY_FONT)],
               line_pts=14, anchor="t")

    c.text(7.56, 5.54, 4.86, 0.42, [(
        "The two push in opposite directions. Both are left out rather than netted "
        "with numbers I cannot source.", 10.5, False, MUTED, BODY_FONT)],
        line_pts=13, anchor="t")

    cov = m["cov"]
    c.callout([
        (f"IDR {ANCHOR_COST / 1e9:.1f}B of cost needs a {m['anchor_pct'] * 100:.1f}% "
         "relative churn reduction on the top-50 pool", 13.5, True, INK, BODY_FONT),
        (f" — and that row is the {cov['ams_total']}-AM programme on the previous "
         f"slide, of which IDR {cov['committed'] / 1e6:,.0f}M is already committed. "
         "A threshold, not a promise.", 13.5, False, INK, BODY_FONT),
    ])

    set_notes(slide, (
        "This is a break-even, not a business case with an uplift in it. I will not put a "
        "fabricated intervention rate in front of you.\n"
        "If asked why the saving is not margin-adjusted: two things are deliberately left out and "
        "they offset -- no gross-margin haircut (makes the bar look easier), and no credit for the "
        "save persisting into year 2 (makes it look harder). I would rather show the omissions than "
        "net them with numbers I cannot source.\n"
        "The bar is RELATIVE, on accounts that are unmanaged today. That is what makes 11.5% a low "
        "bar rather than a heroic one."
    ))


# --------------------------------------------------------------------------
# 4. New slide 12 -- Decision 2's money case.
# --------------------------------------------------------------------------

def build_decision_2(slide, m: dict) -> None:
    c = Canvas(slide)
    u, i = m["unit"], m["illus"]

    c.eyebrow("THE MONEY CASE · DECISION 2")
    c.title("The prepay is a cash decision, not a margin decision")
    c.banner(
        "This does not fund the AM programme.",
        "Decision 2 is priced per rupiah of ARR that takes the offer. Approving one "
        "does not commit you to the other.",
    )

    # -- hero: the working-capital release --------------------------------
    c.card(0.6, 2.3, 4.6, 1.86)
    c.text(0.88, 2.56, 4.04, 0.72,
           [(f"IDR {i['working_capital'] / 1e9:.2f}B", 30, True, NAVY, TITLE_FONT)])
    c.text(0.88, 3.30, 4.04, 0.34,
           [("One-time working capital released", 13, True, INK, BODY_FONT)])
    c.text(0.88, 3.68, 4.04, 0.30, [(
        f"converting {CONVERSION_SHARE:.0%} of the monthly book — IDR "
        f"{m['converted'] / 1e9:.2f}B of ARR", 11, False, MUTED, BODY_FONT)])

    # -- the three assumptions Finance owns -------------------------------
    c.panel(0.6, 4.30, 4.6, 1.78, NAVY, adj=5000)
    c.text(0.86, 4.44, 4.08, 0.26,
           [("THREE ASSUMPTIONS FINANCE OWNS", 10, True, GOLD, BODY_FONT)])
    # Label left, value right, one line each -- wrapped descriptions clipped
    # against the panel's lower border.
    assumptions = [
        ("Cost of capital", f"{COST_OF_CAPITAL:.0%}"),
        ("Value recovered on a failed account", f"{L_RECOVERED:.1f}"),
        ("Haircut on takers' churn rate", f"{ADVERSE_HAIRCUT:.0%}"),
    ]
    for idx, (desc, value) in enumerate(assumptions):
        y = 4.76 + idx * 0.36
        c.text(0.86, y, 2.9, 0.32, [(desc, 10.5, False, ON_DARK_SOFT, BODY_FONT)])
        c.text(3.86, y, 1.08, 0.32,
               [(value, 12, True, GOLD, TITLE_FONT)], align="r")
    c.text(0.86, 5.82, 4.08, 0.24, [(
        "The haircut is adverse selection: the healthiest take it first.",
        9.5, False, ON_DARK_SOFT, BODY_FONT)])

    # -- the rate card ----------------------------------------------------
    c.card(5.44, 2.3, 7.26, 3.78)
    c.text(5.72, 2.49, 6.7, 0.32,
           [("What 8% costs, and what it buys", 17, True, INK, TITLE_FONT)])
    c.text(5.72, 2.88, 6.7, 0.28, [(
        "Every row is a formula, not an estimate — so scale it to whatever "
        "conversion you think we can win.", 11.5, False, MUTED, BODY_FONT)], line_pts=15)

    lx, lw = 5.72, 3.30
    c1x, c2x, cw = 9.12, 10.86, 1.62
    c.text(c1x, 3.30, cw, 0.24, [("PER IDR 1.0B", 9.5, True, MUTED, BODY_FONT)], align="r")
    c.text(c2x, 3.30, cw, 0.24,
           [(f"AT {CONVERSION_SHARE:.0%} OF BOOK", 9.5, True, MUTED, BODY_FONT)], align="r")

    # The first two labels already carry their own parameters (", 8%", ", 11%").
    # The third does not, which is what made it the opaque row -- so it gets its
    # arithmetic on a second line, naming the same three inputs as the
    # assumptions panel opposite. Input rows sit on a 0.40 pitch rather than 0.42
    # to buy that second line its clearance above the totals block.
    retained_formula = (
        f"{W_FORWARD_12M:.1%} churn × {L_RECOVERED} recovered × "
        f"{ADVERSE_HAIRCUT:.0%} haircut"
    )
    rows = [
        (3.58, f"Discount given away, {DISCOUNT:.0%}", "discount", RED, False, None),
        (3.98, f"Carry on {PREPAY_MONTHS_EARLY} months' earlier cash, "
               f"{COST_OF_CAPITAL:.0%}", "carry", TEAL, False, None),
        (4.38, "Revenue kept because the year is locked", "retained", TEAL, False,
         retained_formula),
        (4.88, "Net, per year", "net", INK, True, None),
        (5.30, "One-time cash released", "working_capital", NAVY, True, None),
    ]
    # The two total rows share one highlight block, so they read as a unit.
    c.panel(5.60, 4.84, 6.94, 0.88, PALE_TEAL, adj=5400)
    for y, label, key, colour, bold, sub in rows:
        c.row(y, lx, lw, label, [
            (c1x, cw, money(u[key], "M"), colour),
            (c2x, cw, money(i[key], "M"), colour),
        ], bold=bold, sublabel=sub)

    c.text(5.72, 5.78, 6.7, 0.26, [(
        f"Break-even discount: {m['be_adverse'] * 100:.1f}% under adverse selection, "
        f"{m['be_average'] * 100:.1f}% at average risk — we price at "
        f"{DISCOUNT:.0%}, inside both.", 10.5, True, INK, BODY_FONT)])

    c.callout([
        ("One condition: the prepay must be non-refundable.", 13.5, True, INK, BODY_FONT),
        (f" If it is refundable, the {money(u['retained'], 'M')} retention row goes to "
         "zero, the break-even falls to the carry alone, and we are simply giving away "
         f"{DISCOUNT:.0%}.", 13.5, False, INK, BODY_FONT),
    ])

    set_notes(slide, (
        "Question to expect: does the 8% discount have to be funded by the IDR 1.0B programme "
        "cost? No. Decision 1 is headcount against a covered pool; Decision 2 is priced per "
        "rupiah of ARR that takes the offer. They are independent approvals.\n"
        "Lead with the 4.54B cash release, not the +124M annual net. The net is thin and depends "
        "on three assumptions; the working-capital release does not depend on the churn assumption "
        "at all.\n"
        "The mechanism is already proven -- 1,283 customers bill annually today, so there is no new "
        "billing infrastructure. The opportunity is the 703 wholly-monthly customers, who hold most "
        "of the ARR.\n"
        "If pressed on the 0.6: it is the share of a failed account's annual value we recover by "
        "having been paid up front, from arrears landing roughly 0.6 of the way through the year-2 "
        "window. Finance owns it, and I would take a different number if they have one.\n"
        "HOW TO SAY THE 'REVENUE KEPT' ROW OUT LOUD: on a monthly customer, when they stop paying "
        "we stop collecting -- we lose the rest of the year. If that same year was paid up front and "
        "is non-refundable, we keep it. So the row is the churn we would have suffered, valued at "
        "what being prepaid lets us keep: 12.9% of accounts leave in the next 12 months, we keep "
        "about 0.6 of the year on each one, and I halve it because the healthiest customers take "
        "the offer first. On IDR 1.0B that is 39M. It is not new revenue -- it is revenue we would "
        "otherwise have lost."
    ))


# --------------------------------------------------------------------------
# 5. Slide 9 -- the prepay control line. Slides 2 and 13 -- cross-references.
# --------------------------------------------------------------------------

def patch_slide_9(slide) -> None:
    """Name the commercial offer, its authoriser and its exclusions where the
    renewal sequence already lives -- rather than adding a sixth play card."""
    footer_a = footer_e = card_a = card_e = None
    for sh in slide.shapes:
        if sh.has_text_frame and sh.text_frame.text.startswith("Sequence is diagnose"):
            footer_a = sh
        elif sh.has_text_frame and sh.text_frame.text.startswith("Threshold matters"):
            footer_e = sh
        elif not sh.has_text_frame or not sh.text_frame.text:
            if sh.width == emu(6.02) and sh.height == emu(4.32):
                if sh.left == emu(0.6):
                    card_a = sh
                elif sh.left == emu(7.08):
                    card_e = sh
    assert all((footer_a, footer_e, card_a, card_e)), "slide 9 shapes not found"

    run = footer_a.text_frame.paragraphs[0].runs[-1]
    run.text = run.text.rstrip() + (
        " When one is used it is the 8% non-refundable prepay — authorised by the AM "
        "lead, never to a fleet-event account, never cumulative with a term extension."
    )

    # Three lines of footer need room the 4.32in card does not have. Card E is
    # grown to match even though its text is unchanged, to keep the pair square.
    for card in (card_a, card_e):
        card.height = emu(4.45)
    for footer in (footer_a, footer_e):
        footer.height = emu(0.74)


def patch_the_ask(slide, m: dict) -> None:
    """Decision 1 claimed the top 50 have no owner. Thirteen of them do.

    Also fixes a defect inherited from the source deck: all three decision
    bodies are 8.3in wide starting at 1.76, so they run 0.34in underneath the
    approval pills at 9.72 and the last word of a full line is hidden.
    """
    v = m["cov"]
    bodies, target = [], None
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        if sh.width == emu(8.3) and sh.height == emu(0.66):
            bodies.append(sh)
        if sh.text_frame.text.startswith("They hold IDR 50.9B"):
            target = sh
    assert target is not None, "Decision 1 body not found on the ask slide"
    assert len(bodies) == 3, f"expected 3 decision bodies, found {len(bodies)}"

    # Stop 0.2in short of the approval pills.
    for sh in bodies:
        sh.width = emu(7.76)

    # The 50.9B / 66% framing now lives on the coverage slide immediately
    # before, so this reads shorter and stays on two lines.
    run = target.text_frame.paragraphs[0].runs[0]
    run.text = (
        f"{AMS_TODAY} AMs cover {v['in50_n']} of them today, carrying "
        f"{v['in50_share'] * 100:.0f}% of the expected loss. The other "
        f"{v['unc_n']} have no owner, and none of the 50 has a trigger or a "
        f"renewal calendar. The ask is {AMS_TO_ADD} more AMs — IDR "
        f"{v['incremental'] / 1e6:,.0f}M a year."
    )


def patch_cross_references(prs) -> None:
    """Two insertions shift every later slide by two.

    Searched across the whole deck rather than by slide index, because the
    slides holding these references have themselves moved by the time this runs.
    Source 13 (Corin) lands at 15; source 12 (discipline) lands at 14.
    """
    fixes = [
        ("see slide 13", "see slide 15"),                    # slide 2 -> Corin slide
        ("censoring trap as slide 12", "censoring trap as slide 14"),  # Corin -> discipline
    ]
    for old, new in fixes:
        hits = 0
        for slide in prs.slides:
            for sh in slide.shapes:
                if not sh.has_text_frame:
                    continue
                for para in sh.text_frame.paragraphs:
                    for run in para.runs:
                        if old in run.text:
                            run.text = run.text.replace(old, new)
                            hits += 1
        assert hits == 1, f"expected exactly one {old!r} reference, found {hits}"


def move_slide(prs, from_idx: int, to_idx: int) -> None:
    id_list = prs.slides._sldIdLst
    entries = list(id_list)
    el = entries[from_idx]
    id_list.remove(el)
    id_list.insert(to_idx, el)


# --------------------------------------------------------------------------

def main() -> None:
    print("model:")
    model = build_model()
    check(model)

    work = Path(tempfile.gettempdir()) / "_churn_deck_build.pptx"
    shutil.copy2(SRC_DECK, work)
    prs = Presentation(str(work))
    assert len(prs.slides) == 14, f"expected a 14-slide source deck, got {len(prs.slides)}"

    def add_light_slide():
        """New slide carrying the deck's light background rather than the
        layout's scheme default."""
        s = prs.slides.add_slide(prs.slide_masters[0].slide_layouts[0])
        bg = prs.slides[10]._element.find(qn("p:cSld")).find(qn("p:bg"))
        s._element.find(qn("p:cSld")).insert(0, copy.deepcopy(bg))
        return s

    print("\nbuilding:")
    # Edits to existing slides run first, while source indices still hold.
    build_decision_1(prs.slides[10], model)
    print("  slide 11 rebuilt   -- Decision 1, chart removed")

    patch_the_ask(prs.slides[9], model)
    print("  slide 10 corrected -- 13 of the top 50 already have an owner")

    patch_slide_9(prs.slides[8])
    print("  slide 9 patched    -- prepay authorisation and exclusions")

    build_decision_2(add_light_slide(), model)
    move_slide(prs, len(prs.slides) - 1, 11)
    print("  slide 13 inserted  -- Decision 2, rate card")

    # Coverage goes before the ask, so the correction lands before the request.
    build_coverage_slide(add_light_slide(), model)
    move_slide(prs, len(prs.slides) - 1, 9)
    print("  slide 10 inserted  -- AM coverage and headcount")

    patch_cross_references(prs)
    print("  slides 2, 15       -- forward references renumbered")

    prs.save(str(work))
    shutil.copy2(work, OUT_DECK)
    print(f"\nwrote {OUT_DECK.relative_to(REPO)}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
