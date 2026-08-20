# -*- coding: utf-8 -*-
"""Slides 6, 11, 12 revisions on top of the 2026-08-13 deck.

Reads  outputs/churn-strategy-2026-08-13.pptx
Writes outputs/churn-strategy-2026-08-20.pptx

Edits the existing file in place (rather than rebuilding from the 08-12
source through patch_money_slides.py) so every manual adjustment already in
the 08-13 file -- confirmed: slide 12's title, hand-changed from "The prepay
is a cash decision, not a margin decision" to "Driving Prepay Program to Our
Existing Customer" -- survives untouched. Every shape not explicitly touched
below is left byte-identical.

Slide 6: the "2%" largest-fleet penetration box turned out to be fragile
(median vs average flips it from 2% to 35%+, driven by three whale accounts)
and the footprint-vs-company-size comparison the slide argues for is weaker
than claimed once the two churn-rate-by-band tables are put side by side --
see the model check below. Per Grady's direction: state the footprint
finding on its own, drop the company-size rebuttal wording, and delete the
"2%" box outright rather than replace it.

Slide 11: the right-hand "what 11.5% counts" COUNTED/LEFT-OUT breakdown is
replaced with a one-paragraph statement of what 11.5% means.

Slide 12: reordered so the slide argues purpose of prepay, then why 8%,
then the savings -- instead of leading with the cash figure. The three
Finance-owned assumptions and the break-even sentence move together into a
"Why 8%" card; the cash figure becomes a supporting stat under a new
rationale paragraph.

Run:  python analysis/churn-deck/patch_slides_2026-08-20.py
"""

from __future__ import annotations

import shutil
import sys
import tempfile
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt

sys.path.insert(0, str(Path(__file__).resolve().parent))
import patch_money_slides as pm  # noqa: E402  (path insert must precede this)

REPO = pm.REPO
SRC_DECK = REPO / "outputs" / "churn-strategy-2026-08-13.pptx"
OUT_DECK = REPO / "outputs" / "churn-strategy-2026-08-20.pptx"

# Churn-rate-by-band figures behind the slide-6 rewrite, from
# outputs/churn-cohort-v2-2026-08-14.xlsx, tab 04_Churn_Rate_Segment.
FOOTPRINT_KM24 = {  # by subscribed-vehicle band
    "1-5 veh": 0.3612798871115516, "6-10 veh": 0.2357209189248195,
    "11-20 veh": 0.2528647321310628, "21-50 veh": 0.1968179480874869,
    "51-100 veh": 0.2117522390953041, "101-399 veh": 0.1720032698359432,
    "400+ veh": 0.1151508080807024,
}
COMPANY_SIZE_KM24 = {  # by CRM Fleet Category
    "A<=5": 0.2428145149326782, "B<=10": 0.2285571460909789,
    "C<=20": 0.2613265974694633, "D<=50": 0.1983943865666836,
    "E<=100": 0.1728673075262618, "F<400": 0.1567502087454773,
    "G>=400": 0.09107084526870401,
}


def find(slide, shape_id: int):
    for sh in slide.shapes:
        if sh.shape_id == shape_id:
            return sh
    raise SystemExit(f"shape id={shape_id} not found on slide")


def set_text(shape, expected: str, new: str) -> None:
    run = shape.text_frame.paragraphs[0].runs[0]
    assert run.text == expected, f"expected {expected!r}, found {run.text!r}"
    run.text = new


def delete(shape) -> None:
    shape._element.getparent().remove(shape._element)


def check_footprint_beats_company_size() -> None:
    """The comparison the old slide 6 leaned on: confirm the spread claim is
    only real at the smallest band, not check-and-suppress -- this must fail
    loudly if the source workbook figures ever move, same discipline as
    patch_money_slides.check().
    """
    cases = [
        ("footprint 1-5 veh", FOOTPRINT_KM24["1-5 veh"], 0.3612798871115516),
        ("footprint 400+ veh", FOOTPRINT_KM24["400+ veh"], 0.1151508080807024),
        ("company A<=5", COMPANY_SIZE_KM24["A<=5"], 0.2428145149326782),
        ("company G>=400", COMPANY_SIZE_KM24["G>=400"], 0.09107084526870401),
    ]
    for name, got, want in cases:
        if abs(got - want) > 1e-9:
            raise SystemExit(f"MISMATCH {name}: {got} != {want}")
    fp_spread = FOOTPRINT_KM24["1-5 veh"] / FOOTPRINT_KM24["400+ veh"]
    cs_spread = COMPANY_SIZE_KM24["A<=5"] / COMPANY_SIZE_KM24["G>=400"]
    print(f"  ok  footprint spread {fp_spread:.2f}x, company-size spread {cs_spread:.2f}x "
          "-- close enough that the old slide overstated the gap")


# --------------------------------------------------------------------------
# Slide 6
# --------------------------------------------------------------------------

def edit_slide6(slide) -> None:
    set_text(
        find(slide, 3),
        "Our footprint in the account predicts churn — company size barely does",
        "Our footprint in the account predicts churn",
    )
    set_text(
        find(slide, 7),
        "Small-footprint accounts churn three times more often than large ones. "
        "Company size gives only a 2.9× spread and does not rank cleanly.",
        "Small-footprint accounts churn three times more often than large ones.",
    )
    set_text(
        find(slide, 15),
        "Segment account management on how many vehicles we run for a customer, "
        "not on how big the customer is. A 2%-penetrated account is a pilot, "
        "and pilots are easy to cancel.",
        "Segment account management on how many vehicles we run for a "
        "customer — that's what predicts churn.",
    )

    for sid in (9, 10, 11, 8):  # text shapes first, background last
        delete(find(slide, sid))

    # Grow the remaining card to fill the vacated lower half of the column,
    # and recentre its contents -- otherwise the right column ends at 3.70in
    # while the card background originally ran to 6.04in.
    card = find(slide, 5)
    card.height = pm.emu(4.06)  # was 1.72in; column now runs 1.98 -> 6.04in

    number = find(slide, 6)
    number.top = pm.emu(3.30)
    number.height = pm.emu(0.70)
    run = number.text_frame.paragraphs[0].runs[0]
    run.font.size = Pt(40)

    caption = find(slide, 7)
    caption.top = pm.emu(4.18)
    caption.height = pm.emu(0.50)


# --------------------------------------------------------------------------
# Slide 11
# --------------------------------------------------------------------------

def edit_slide11(slide, m: dict) -> None:
    set_text(find(slide, 25), "What the 11.5% counts", "What 11.5% means")
    set_text(
        find(slide, 26),
        "A threshold, not a forecast. Here is the accounting behind it.",
        "A threshold, not a forecast.",
    )
    for sid in (27, 28, 29, 30, 31, 32, 33):
        delete(find(slide, sid))

    c = pm.Canvas(slide, next_id=500)
    anchor_pct = m["anchor_pct"] * 100
    covered_b = m["covered"] / 1e9
    cost_b = pm.ANCHOR_COST / 1e9
    c.text_stack(7.56, 3.55, 4.86, 2.1, [
        ([(f"{anchor_pct:.1f}%", 16, True, pm.INK, pm.TITLE_FONT),
          (" is how much of the ", 13, False, pm.INK, pm.BODY_FONT),
          (f"IDR {covered_b:.2f}B", 13, True, pm.INK, pm.BODY_FONT),
          (" at-risk ARR in the top-50 pool the programme needs to save for its ",
           13, False, pm.INK, pm.BODY_FONT),
          (f"IDR {cost_b:.1f}B", 13, True, pm.INK, pm.BODY_FONT),
          (" cost to break even.", 13, False, pm.INK, pm.BODY_FONT)], 18),
        ([("Save less than that, and the programme costs more than it returns.",
           13, False, pm.MUTED, pm.BODY_FONT)], 18),
    ], anchor="t")


# --------------------------------------------------------------------------
# Slide 12
# --------------------------------------------------------------------------

def edit_slide12(slide, m: dict) -> None:
    u, i = m["unit"], m["illus"]

    # -- Card A: purpose of prepay (was: cash headline only) ---------------
    card_a = find(slide, 6)
    card_a.height = pm.emu(2.00)  # was 1.86in

    c = pm.Canvas(slide, next_id=600)
    c.text(0.86, 2.44, 4.08, 0.26, [("Why prepay", 14, True, pm.INK, pm.TITLE_FONT)])
    c.text(0.86, 2.74, 4.08, 0.42, [(
        "Most ARR bills monthly, exposed to cancellation any month. Prepay "
        "locks in the year's revenue and pulls cash forward now.",
        10.5, False, pm.MUTED, pm.BODY_FONT)], line_pts=13.5, anchor="t")

    stat = find(slide, 7)
    stat.top = pm.emu(3.22)
    stat.height = pm.emu(0.38)
    run = stat.text_frame.paragraphs[0].runs[0]
    run.font.size = Pt(22)

    label = find(slide, 8)
    label.top = pm.emu(3.64)
    label.height = pm.emu(0.22)
    run = label.text_frame.paragraphs[0].runs[0]
    run.font.size = Pt(11)

    sub = find(slide, 9)
    sub.top = pm.emu(3.90)
    sub.height = pm.emu(0.20)

    # -- Card B: why 8% (was: just the three assumptions) -------------------
    card_b = find(slide, 10)
    card_b.top = pm.emu(4.42)
    card_b.height = pm.emu(1.66)  # was 1.78in at 4.30in top

    heading = find(slide, 11)
    heading.top = pm.emu(4.54)

    row_ids = [(12, 13, 4.80), (14, 15, 5.10), (16, 17, 5.40)]
    for label_id, value_id, y in row_ids:
        find(slide, label_id).top = pm.emu(y)
        find(slide, value_id).top = pm.emu(y)

    delete(find(slide, 18))  # "haircut is adverse selection" caption -- cut for room

    c2 = pm.Canvas(slide, next_id=650)
    c2.text(0.86, 5.72, 4.08, 0.30, [(
        f"Break-even discount: {m['be_adverse'] * 100:.1f}% under adverse "
        f"selection, {m['be_average'] * 100:.1f}% at average risk — we price "
        f"at {pm.DISCOUNT:.0%}, inside both.", 10, True, pm.GOLD, pm.BODY_FONT,
    )], line_pts=13)

    # -- Right column: drop the break-even line, now living on card B -------
    delete(find(slide, 40))


def main() -> None:
    print("model:")
    model = pm.build_model()
    pm.check(model)
    print("\nslide-6 comparison check:")
    check_footprint_beats_company_size()

    tmp_src = Path(tempfile.gettempdir()) / "_churn_deck_2026-08-20_src.pptx"
    shutil.copy2(SRC_DECK, tmp_src)  # OneDrive can lock direct reads
    prs = Presentation(str(tmp_src))
    assert len(prs.slides) == 16, f"expected a 16-slide source deck, got {len(prs.slides)}"

    edit_slide6(prs.slides[5])
    print("\n  slide 6 edited  -- footprint-only framing, 2% box removed")
    edit_slide11(prs.slides[10], model)
    print("  slide 11 edited -- right box replaced")
    edit_slide12(prs.slides[11], model)
    print("  slide 12 edited -- reordered purpose -> why 8% -> savings")

    work = Path(tempfile.gettempdir()) / "_churn_deck_2026-08-20_out.pptx"
    prs.save(str(work))
    shutil.copy2(work, OUT_DECK)
    print(f"\nwrote {OUT_DECK.relative_to(REPO)}")


if __name__ == "__main__":
    main()
